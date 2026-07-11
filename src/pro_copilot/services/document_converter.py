import logging
import shutil
from datetime import datetime, timezone
from pathlib import Path

from pro_copilot.config import settings

logger = logging.getLogger(__name__)


async def run_document_conversion() -> None:
    """掃描 raw_logs/incoming/ 目錄，將裡面的 office/pdf 檔案轉換為 Markdown 存入 raw_logs/documents/。"""
    incoming_dir: Path = settings.raw_logs_dir / "incoming"
    archive_dir: Path = incoming_dir / "archive"
    output_dir: Path = settings.raw_logs_dir / "documents"

    # 確保資料夾存在
    incoming_dir.mkdir(parents=True, exist_ok=True)
    archive_dir.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    supported_extensions = {".pdf", ".docx", ".pptx", ".xlsx"}
    
    # 篩選出待轉換的檔案
    files_to_convert = [
        f for f in incoming_dir.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]

    if not files_to_convert:
        logger.info("沒有找到待轉換的原始文件。")
        return

    logger.info("找到 %d 個待轉換檔案，開始進行轉換...", len(files_to_convert))

    for file_path in files_to_convert:
        logger.info("正在轉換檔案: %s", file_path.name)
        ext = file_path.suffix.lower()
        try:
            if ext == ".pdf":
                text = _parse_pdf(file_path)
            elif ext == ".docx":
                text = _parse_docx(file_path)
            elif ext == ".pptx":
                text = _parse_pptx(file_path)
            elif ext == ".xlsx":
                text = _parse_xlsx(file_path)
            else:
                continue

            # 組合 metadata frontmatter
            now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
            md_content = "\n".join([
                "---",
                f"date: {now_str}",
                "source: document",
                f"original_file: {file_path.name}",
                "---",
                "",
                text
            ])

            # 寫出 Markdown 檔案
            output_file = output_dir / f"{file_path.stem}.md"
            output_file.write_text(md_content, encoding="utf-8")
            logger.info("檔案 %s 轉換成功 -> %s", file_path.name, output_file.name)

            # 將原檔案移入封存區，避免重複處理
            dest_archive_path = archive_dir / file_path.name
            # 如果封存區已有同名檔案，加上時間戳避免衝突
            if dest_archive_path.exists():
                timestamp = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S")
                dest_archive_path = archive_dir / f"{file_path.stem}_{timestamp}{ext}"

            shutil.move(str(file_path), str(dest_archive_path))
            logger.info("已將原始檔案移入封存目錄: %s", dest_archive_path.name)

        except Exception as e:
            logger.error("轉換檔案 %s 失敗: %s", file_path.name, e, exc_info=True)


def _clean_spaced_out_text(text: str) -> str:
    """清理 PDF 提取文字時常見的字元間隔多餘空格問題。"""
    import re
    cleaned_lines = []
    for line in text.splitlines():
        # 以多於一個空格（兩個或以上）來切分單字/片語組
        words = re.split(r' {2,}', line.strip())
        cleaned_words = []
        for word in words:
            non_space_chars = [c for c in word if c != ' ']
            if not non_space_chars:
                continue
            # 若單字長度大於 1 且剛好是 (非空格字元數 * 2 - 1)，檢查是否為交替空格字元（如 "P y t h o n"）
            if len(word) > 1 and len(word) == 2 * len(non_space_chars) - 1:
                is_spaced_out = True
                for idx, char in enumerate(word):
                    if idx % 2 == 0:
                        if char == ' ':
                            is_spaced_out = False
                            break
                    else:
                        if char != ' ':
                            is_spaced_out = False
                            break
                if is_spaced_out:
                    cleaned_words.append("".join(non_space_chars))
                else:
                    cleaned_words.append(word)
            else:
                cleaned_words.append(word)
        cleaned_lines.append(" ".join(cleaned_words))
    return "\n".join(cleaned_lines)


def _parse_pdf(file_path: Path) -> str:
    """解析 PDF 檔案提取純文字。"""
    from pypdf import PdfReader
    
    reader = PdfReader(file_path)
    text_parts = []
    for i, page in enumerate(reader.pages, 1):
        page_text = page.extract_text()
        if page_text and page_text.strip():
            cleaned_text = _clean_spaced_out_text(page_text)
            text_parts.append(f"## Page {i}\n")
            text_parts.append(cleaned_text.strip())
            text_parts.append("")
            
    return "\n".join(text_parts)


def _parse_docx(file_path: Path) -> str:
    """解析 Word (.docx) 檔案提取段落與表格文字。"""
    import docx
    
    doc = docx.Document(file_path)
    text_parts = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
            
    if doc.tables:
        text_parts.append("\n## 表格資料")
        for table in doc.tables:
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " ").replace("\r", "") for cell in row.cells]
                text_parts.append("| " + " | ".join(cells) + " |")
                if r_idx == 0:
                    text_parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
            text_parts.append("")
            
    return "\n".join(text_parts)


def _parse_pptx(file_path: Path) -> str:
    """解析 PowerPoint (.pptx) 檔案提取投影片文字與表格。"""
    from pptx import Presentation
    
    prs = Presentation(file_path)
    text_parts = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"## Slide {i}")
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            if shape.has_table:
                table = shape.table
                table_lines = []
                for r_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ").replace("\r", "") for cell in row.cells]
                    table_lines.append("| " + " | ".join(cells) + " |")
                    if r_idx == 0:
                        table_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                slide_texts.append("\n" + "\n".join(table_lines) + "\n")
        
        if slide_texts:
            text_parts.extend(slide_texts)
            text_parts.append("")
            
    return "\n".join(text_parts)


def _parse_xlsx(file_path: Path) -> str:
    """解析 Excel (.xlsx) 檔案，將工作表內容轉換為 Markdown 表格。"""
    import openpyxl
    
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text_parts = []
    
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
            
        # 移除尾部的全空行
        while rows and not any(cell is not None and str(cell).strip() != "" for cell in rows[-1]):
            rows.pop()
            
        if not rows:
            continue
            
        text_parts.append(f"## Sheet: {sheet.title}\n")
        
        for r_idx, row in enumerate(rows):
            formatted_cells = []
            for cell in row:
                if cell is None:
                    formatted_cells.append("")
                else:
                    # 取代換行以避免破壞 Markdown 表格語法
                    formatted_cells.append(str(cell).replace("\n", " ").replace("\r", ""))
            
            text_parts.append("| " + " | ".join(formatted_cells) + " |")
            if r_idx == 0:
                text_parts.append("| " + " | ".join(["---"] * len(row)) + " |")
                
        text_parts.append("")
        
    return "\n".join(text_parts)
